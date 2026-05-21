
from __future__ import annotations

import io
import math
import re
import pandas as pd
import string
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from io import BytesIO

try:
    import pdfplumber
    _PDF_OK = True
except ImportError:
    _PDF_OK = False

try:
    from docx import Document as DocxDocument
    _DOCX_OK = True
except ImportError:
    _DOCX_OK = False

try:
    from pptx import Presentation
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False
 
SENT_SPLIT_RE = re.compile(r"(?<=[.!?])\s+")

_STOP_WORDS = {
    "a","about","above","after","again","against","all","also","am","an","and",
    "any","are","aren't","as","at","be","because","been","before","being",
    "below","between","both","but","by","can","couldn't","did","didn't","do",
    "does","doesn't","doing","don","don't","down","during","each","few","for",
    "from","further","get","had","hadn't","has","hasn't","have","haven't",
    "having","he","he'd","he'll","he's","her","here","here's","hers","herself",
    "him","himself","his","how","how's","i","i'd","i'll","i'm","i've","if","in",
    "into","is","isn't","it","it's","its","itself","let","let's","me","more",
    "most","mustn't","my","myself","no","nor","not","of","off","on","once","only",
    "or","other","ought","our","ours","ourselves","out","over","own","same",
    "shan't","she","she'd","she'll","she's","should","shouldn't","so","some",
    "such","than","that","that's","the","their","theirs","them","themselves",
    "then","there","there's","therefore","these","they","they'd","they'll",
    "they're","they've","this","those","through","to","too","under","until","up",
    "us","very","was","wasn't","we","we'd","we'll","we're","we've","were",
    "weren't","what","what's","when","when's","where","where's","which","while",
    "who","who's","whom","why","why's","will","with","won't","would","wouldn't",
    "you","you'd","you'll","you're","you've","your","yours","yourself","yourselves",
    # domain stop
    "said","say","says","also","however","therefore","thus","hence","one","two",
    "three","four","five","six","seven","eight","nine","ten","per","within",
    "between","across","among","towards","regarding","including","following",
    "based","may","shall","will","can","could","use","used","using","make",
    "made","need","needs","via","many","much","several","various",
}

_SAP_TERMS = {
    "fi","co","sd","mm","pp","qm","hr","fico","fiori","s4hana","ecc","erp",
    "abap","bapi","idoc","bdc","sap","hana","ariba","successfactors","concur",
    "po","gr","ir","ap","ar","gl","am","wbs","profit center","cost center",
    "plant","storage location","material","vendor","customer","invoice",
    "purchase order","goods receipt","invoice verification","payment run",
    "clearing","dunning","valuation","depreciation","asset","posting","period",
    "fiscal year","company code","controlling area","operating concern",
}

def _tokenize(text: str) -> List[str]:
    text = text.lower()
    text = re.sub(r"[^a-z0-9\s]", " ", text)
    tokens = text.split()
    return [t for t in tokens if t not in _STOP_WORDS and len(t) > 1]

def _sentences(text: str) -> List[str]:
    # split on sentence boundaries while keeping context
    raw = re.split(r"(?<=[.!?])\s+(?=[A-Z\"'])|(?<=\n)\s*(?=[A-Z\"'])", text)
    out = []
    for s in raw:
        s = s.strip()
        if len(s.split()) >= 5:  # ignore very short fragments
            out.append(s)
    return out

def _cosine(v1: Dict[str, float], v2: Dict[str, float]) -> float:
    keys = set(v1) & set(v2)
    if not keys:
        return 0.0
    dot = sum(v1[k] * v2[k] for k in keys)
    mag1 = math.sqrt(sum(x * x for x in v1.values()))
    mag2 = math.sqrt(sum(x * x for x in v2.values()))
    if mag1 == 0 or mag2 == 0:
        return 0.0
    return dot / (mag1 * mag2)

def _extract_text_pdf(data: bytes) -> str:
    if not _PDF_OK:
        return ""
    try:
        pages = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                pages.append(t)
        return "\n".join(pages)
    except Exception:
        return ""

def _extract_text_docx(data: bytes) -> str:
    if not _DOCX_OK:
        return ""
    try:
        doc = DocxDocument(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""

def _extract_text_pptx(data: bytes) -> str:
    if not _PPTX_OK:
        return ""
    try:
        prs = Presentation(io.BytesIO(data))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
        return "\n".join(lines)
    except Exception:
        return ""

def _extract_text(filename: str, data: bytes) -> str:
    ext = Path(filename).suffix.lower()

    # 1) Rich text formats
    if ext == ".pdf":
        return _extract_text_pdf(data)
    if ext in (".docx", ".doc"):
        return _extract_text_docx(data)
    if ext in (".pptx", ".ppt"):
        return _extract_text_pptx(data)

    # 2) CSV invoice exports → convert rows to sentences
    if ext == ".csv":
        try:
            df = pd.read_csv(BytesIO(data))

            # try to use common invoice columns if they exist
            vendor_col = "Vendor" if "Vendor" in df.columns else None
            inv_col    = "InvoiceNumber" if "InvoiceNumber" in df.columns else None
            date_col   = "InvoiceDate" if "InvoiceDate" in df.columns else None
            amt_col    = "Amount" if "Amount" in df.columns else None
            curr_col   = "Currency" if "Currency" in df.columns else None
            status_col = "PaymentStatus" if "PaymentStatus" in df.columns else None

            lines = []
            for _, row in df.iterrows():
                vendor = str(row.get(vendor_col, "")) if vendor_col else ""
                inv    = str(row.get(inv_col, "")) if inv_col else ""
                date   = str(row.get(date_col, "")) if date_col else ""
                amt    = row.get(amt_col, "")
                curr   = str(row.get(curr_col, "")) if curr_col else ""
                status = str(row.get(status_col, "")) if status_col else ""

                # build one readable sentence per invoice
                sentence = f"Invoice {inv} from {vendor} dated {date} amount {amt} {curr} has status {status}."
                lines.append(sentence)

            return "\n".join(lines)
        except Exception:
            return ""

    if ext in (".txt", ".md", ".log"):
        try:
            return data.decode("utf-8", errors="replace")
        except Exception:
            return ""
        
    if ext in (".xlsx", ".xls"):
        try:
            df = pd.read_excel(BytesIO(data))
            # reuse the same logic as CSV
            vendor_col = "Vendor" if "Vendor" in df.columns else None
            inv_col    = "InvoiceNumber" if "InvoiceNumber" in df.columns else None
            date_col   = "InvoiceDate" if "InvoiceDate" in df.columns else None
            amt_col    = "Amount" if "Amount" in df.columns else None
            curr_col   = "Currency" if "Currency" in df.columns else None
            status_col = "PaymentStatus" if "PaymentStatus" in df.columns else None

            lines = []
            for _, row in df.iterrows():
                vendor = str(row.get(vendor_col, "")) if vendor_col else ""
                inv    = str(row.get(inv_col, "")) if inv_col else ""
                date   = str(row.get(date_col, "")) if date_col else ""
                amt    = row.get(amt_col, "")
                curr   = str(row.get(curr_col, "")) if curr_col else ""
                status = str(row.get(status_col, "")) if status_col else ""

                sentence = f"Invoice {inv} from {vendor} dated {date} amount {amt} {curr} has status {status}."
                lines.append(sentence)

            return "\n".join(lines)
        except Exception:
            return ""

    try:
        return data.decode("utf-8", errors="replace")
    except Exception:
        return ""

class _TFIDF:
    """Lightweight TF-IDF over a list of sentences."""

    def __init__(self, sentences: List[str]):
        self._docs: List[List[str]] = [_tokenize(s) for s in sentences]
        n = len(self._docs)
        # document frequency
        df: Counter = Counter()
        for tokens in self._docs:
            df.update(set(tokens))
        # IDF with smoothing
        self._idf: Dict[str, float] = {
            w: math.log((n + 1) / (df[w] + 1)) + 1.0
            for w in df
        }

    def vector(self, idx: int) -> Dict[str, float]:
        tokens = self._docs[idx]
        if not tokens:
            return {}
        tf = Counter(tokens)
        total = len(tokens)
        return {w: (tf[w] / total) * self._idf.get(w, 1.0) for w in tf}

    def corpus_idf(self) -> Dict[str, float]:
        return self._idf

def _textrank(
    vectors: List[Dict[str, float]],
    damping: float = 0.85,
    max_iter: int = 60,
    tol: float = 1e-4,
) -> List[float]:
    n = len(vectors)
    if n == 0:
        return []

    # build similarity matrix
    sim: List[List[float]] = [[0.0] * n for _ in range(n)]
    for i in range(n):
        for j in range(i + 1, n):
            s = _cosine(vectors[i], vectors[j])
            sim[i][j] = s
            sim[j][i] = s

    # row-normalise
    for i in range(n):
        row_sum = sum(sim[i])
        if row_sum > 0:
            sim[i] = [v / row_sum for v in sim[i]]

    # power iteration
    scores = [1.0 / n] * n
    for _ in range(max_iter):
        new_scores = [
            (1 - damping) / n + damping * sum(sim[j][i] * scores[j] for j in range(n))
            for i in range(n)
        ]
        if max(abs(new_scores[i] - scores[i]) for i in range(n)) < tol:
            break
        scores = new_scores

    return scores

def _extract_keywords(idf: Dict[str, float], all_tokens: List[str], top_n: int = 15) -> List[str]:
    """TF-IDF based keyword ranking over the full corpus."""
    total = len(all_tokens)
    if total == 0:
        return []
    tf = Counter(all_tokens)
    scores = {w: (tf[w] / total) * idf.get(w, 1.0) for w in tf}
    return [w for w, _ in sorted(scores.items(), key=lambda x: -x[1])[:top_n]]

def _rake_phrases(text: str, top_n: int = 10) -> List[str]:
    """Rapid Automatic Keyword Extraction (RAKE-lite)."""
    sentences = re.split(r"[.!?,;\n]", text)
    stop_pat = r"\b(?:" + "|".join(re.escape(w) for w in sorted(_STOP_WORDS, key=len, reverse=True)) + r")\b"
    phrase_list: List[str] = []
    for s in sentences:
        parts = re.split(stop_pat, s.lower())
        for part in parts:
            phrase = re.sub(r"[^a-z0-9 ]", "", part).strip()
            if 2 <= len(phrase.split()) <= 5 and len(phrase) > 4:
                phrase_list.append(phrase)

    # score by word frequency / degree
    word_freq: Counter = Counter()
    word_degree: Counter = Counter()
    for phrase in phrase_list:
        words = phrase.split()
        for w in words:
            word_freq[w] += 1
            word_degree[w] += len(words)

    phrase_scores: Dict[str, float] = {}
    for phrase in set(phrase_list):
        words = phrase.split()
        score = sum((word_degree[w] + word_freq[w]) / max(word_freq[w], 1) for w in words)
        phrase_scores[phrase] = score

    ranked = sorted(phrase_scores, key=lambda x: -phrase_scores[x])
    # deduplicate subsets
    final: List[str] = []
    for phrase in ranked:
        if not any(phrase in existing for existing in final):
            final.append(phrase)
        if len(final) >= top_n:
            break
    return final

_CURRENCY_PAT = re.compile(
    r"(?:USD|EUR|GBP|INR|₹|\$|€|£)\s?[\d,]+(?:\.\d{1,2})?(?:\s?(?:million|billion|lakh|crore|thousand|k|M|B))?",
    re.IGNORECASE,
)
_PERCENT_PAT = re.compile(r"\d+(?:\.\d+)?\s?%")
_DATE_PAT = re.compile(
    r"\b(?:\d{1,2}[\/\-]\d{1,2}[\/\-]\d{2,4}|"
    r"(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s+\d{1,2},?\s+\d{4}|"
    r"\d{4}[\/\-]\d{2}[\/\-]\d{2})\b",
    re.IGNORECASE,
)
_SAP_CODE_PAT = re.compile(r"\b(?:FB[0-9LVZ]+|VK\d+|ME\d+|MIRO|MIGO|F110|FK\d+|XK\d+)\b")
_ORG_SUFFIX_PAT = re.compile(
    r"\b[A-Z][a-z]+(?: [A-Z][a-z]+)*\s+(?:Ltd|LLC|GmbH|AG|Corp|Inc|Pvt|Private|Limited|Group|Solutions|Services|Technologies|Systems)\b"
)

def _extract_entities(text: str) -> Dict[str, List[str]]:
    entities: Dict[str, List[str]] = {
        "currency_figures": [],
        "percentages": [],
        "dates": [],
        "sap_codes": [],
        "organisations": [],
    }
    entities["currency_figures"] = list(dict.fromkeys(_CURRENCY_PAT.findall(text)))[:8]
    entities["percentages"] = list(dict.fromkeys(_PERCENT_PAT.findall(text)))[:8]
    entities["dates"] = list(dict.fromkeys(_DATE_PAT.findall(text)))[:8]
    entities["sap_codes"] = list(dict.fromkeys(_SAP_CODE_PAT.findall(text)))[:8]
    entities["organisations"] = list(dict.fromkeys(_ORG_SUFFIX_PAT.findall(text)))[:8]
    return entities


def _is_heading(sentence: str) -> bool:
    s = sentence.strip()
    if len(s) > 120:
        return False
    if s.isupper() and len(s) > 3:
        return True
    if re.match(r"^\d+[\.\)]\s+[A-Z]", s):
        return True
    if re.match(r"^[A-Z][^.!?]{3,60}$", s):
        return True
    return False


def _classify_doc_type(text: str, filename: str) -> str:
    tl = text.lower()
    name = filename.lower()
    if any(w in tl for w in ["invoice", "billing", "payment", "vendor", "accounts payable"]):
        return "SAP Finance / Invoice Document"
    if any(w in tl for w in ["purchase order", "procurement", "material", "goods receipt"]):
        return "SAP Procurement Document"
    if any(w in tl for w in ["sales order", "customer", "delivery", "shipment", "revenue"]):
        return "SAP Sales Document"
    if any(w in tl for w in ["policy", "procedure", "guideline", "compliance", "audit"]):
        return "Policy / Compliance Document"
    if any(w in tl for w in ["contract", "agreement", "terms", "clause", "obligations"]):
        return "Contract / Legal Document"
    if any(w in tl for w in ["report", "analysis", "data", "performance", "kpi", "metric"]):
        return "Business Report / Analysis"
    if any(w in tl for w in ["presentation", "slide", "overview", "agenda"]):
        return "Presentation"
    return "General SAP Business Document"

_POSITIVE_WORDS = {
    "increase","growth","improve","improved","success","positive","benefit",
    "gain","profit","efficient","effective","achievement","strong","excellent",
    "compliant","approved","resolved","optimised","savings","reduced","better",
}
_NEGATIVE_WORDS = {
    "risk","fraud","error","violation","discrepancy","missing","overdue","penalty",
    "dispute","reject","failed","delay","non-compliant","breach","loss","concern",
    "anomaly","suspicious","exception","blocked","cancelled","reversed","incorrect",
}

def _sentiment(text: str) -> str:
    tokens = set(text.lower().split())
    pos = len(tokens & _POSITIVE_WORDS)
    neg = len(tokens & _NEGATIVE_WORDS)
    if neg > pos + 1:
        return "⚠️ Risk / Concern Signals Detected"
    if pos > neg + 1:
        return "✅ Positive / Compliant Tone"
    return "🔵 Neutral / Informational"


def _reading_time(word_count: int) -> str:
    mins = max(1, round(word_count / 200))
    return f"{mins} min read"


def summarise_documents(file_pairs, num_points):
    """
    file_pairs: list of (filename, bytes) from Flask
    num_points: int, how many key sentences to return in combined summary
    """
    per_doc = []
    all_sentences = []
    total_words = 0
    total_sentences = 0
    errors = []

    for filename, raw in file_pairs:
        try:
            # turn bytes into text
            text = _extract_text(filename, raw)
            text = text.replace("\r", " ").replace("\n", " ").strip()
            if not text:
                per_doc.append({"filename": filename, "word_count": 0, "key_points": []})
                continue

            # split into sentences (simple splitter)
            sentences = [s.strip() for s in SENT_SPLIT_RE.split(text) if s.strip()]
            words = text.split()

            total_words += len(words)
            total_sentences += len(sentences)

            # simple per-doc highlights: first 3 sentences
            per_doc.append({
                "filename": filename,
                "word_count": len(words),
                "key_points": sentences[:3]
            })

            all_sentences.extend(sentences)

        except Exception as e:
            errors.append({"filename": filename, "error": str(e)})

    # combined summary: first num_points sentences across all docs
    summary_sentences = all_sentences[:num_points]
    summary_text = " ".join(summary_sentences)

    stats = {
        "total_documents": len(per_doc),
        "total_words": total_words,
        "total_sentences_analysed": total_sentences,
        "sentences_in_summary": len(summary_sentences),
    }

    return {
        "stats": stats,
        "per_doc": per_doc,
        "summary": summary_text,
        "bullet_points": summary_sentences,
        "errors": errors,
    }



def _detect_sections(
    sentences: List[str], scores: List[float]
) -> List[Dict[str, Any]]:
    """
    Group consecutive sentences into logical sections using heading detection.
    Return top-2 sentences per section.
    """
    sections: List[Dict[str, Any]] = []
    current_title = "Introduction"
    current_sents: List[Tuple[int, str]] = []

    def _flush(title: str, sents: List[Tuple[int, str]]) -> None:
        if not sents:
            return
        # pick top 2 by score
        ranked = sorted(sents, key=lambda x: -scores[x[0]])
        top_sents = [s for _, s in sorted(ranked[:2], key=lambda x: x[0])]
        sections.append({"section": title, "key_sentences": top_sents})

    for i, sent in enumerate(sentences):
        if _is_heading(sent):
            _flush(current_title, current_sents)
            current_title = sent.strip()
            current_sents = []
        else:
            current_sents.append((i, sent))

    _flush(current_title, current_sents)
    return sections[:12]  # cap at 12 sections