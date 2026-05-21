from io import BytesIO
import re
import pdfplumber
from docx import Document
from pptx import Presentation


def clean_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\xa0", " ")
    text = text.replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def extract_pdf(file_bytes, file_name):
    docs = []
    with pdfplumber.open(BytesIO(file_bytes)) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            text = clean_text(text)
            if text:
                docs.append({
                    "text": text,
                    "source": {
                        "file_name": file_name,
                        "location": f"Page {page_num}",
                        "doc_type": "pdf"
                    }
                })
    return docs


def extract_docx(file_bytes, file_name):
    doc = Document(BytesIO(file_bytes))
    docs = []

    para_buffer = []
    block_num = 1

    for p in doc.paragraphs:
        txt = clean_text(p.text)
        if txt:
            para_buffer.append(txt)

        if len(" ".join(para_buffer)) > 1200:
            docs.append({
                "text": "\n".join(para_buffer),
                "source": {
                    "file_name": file_name,
                    "location": f"Section {block_num}",
                    "doc_type": "docx"
                }
            })
            para_buffer = []
            block_num += 1

    if para_buffer:
        docs.append({
            "text": "\n".join(para_buffer),
            "source": {
                "file_name": file_name,
                "location": f"Section {block_num}",
                "doc_type": "docx"
            }
        })

    return docs


def extract_pptx(file_bytes, file_name):
    prs = Presentation(BytesIO(file_bytes))
    docs = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = clean_text(shape.text)
                if txt:
                    texts.append(txt)

        slide_text = "\n".join(texts).strip()
        if slide_text:
            docs.append({
                "text": slide_text,
                "source": {
                    "file_name": file_name,
                    "location": f"Slide {slide_num}",
                    "doc_type": "pptx"
                }
            })

    return docs


def extract_file(uploaded_file):
    file_name = uploaded_file.filename
    ext = file_name.lower().split(".")[-1]
    file_bytes = uploaded_file.read()

    if ext == "pdf":
        return extract_pdf(file_bytes, file_name)
    if ext == "docx":
        return extract_docx(file_bytes, file_name)
    if ext == "pptx":
        return extract_pptx(file_bytes, file_name)

    return []